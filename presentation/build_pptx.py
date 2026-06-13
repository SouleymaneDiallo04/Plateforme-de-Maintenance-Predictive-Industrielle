# -*- coding: utf-8 -*-
"""Generateur du PowerPoint de soutenance PrognoSense (15 slides + notes)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Palette (couleurs de l'application) ----
TEAL   = RGBColor(0, 130, 160)
DARK   = RGBColor(15, 42, 56)
AMBER  = RGBColor(222, 150, 30)
GREEN  = RGBColor(0, 150, 100)
GREY   = RGBColor(80, 90, 100)
LIGHT  = RGBColor(244, 248, 250)
WHITE  = RGBColor(255, 255, 255)
INK    = RGBColor(30, 41, 51)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tf

def setrun(r, text, size, color, bold=False, italic=False, font="Calibri"):
    r.text = text; r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font

def footer(s, n):
    rect(s, Inches(0), Inches(7.15), SW, Pt(2), LIGHT)
    tf = textbox(s, Inches(0.5), Inches(7.05), Inches(9), Inches(0.4))
    setrun(tf.paragraphs[0].add_run(),
           "PrognoSense — Plateforme de Maintenance Prédictive Industrielle", 9, GREY)
    tf2 = textbox(s, Inches(11.3), Inches(7.05), Inches(1.6), Inches(0.4))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    setrun(p.add_run(), f"{n} / 15", 9, TEAL, bold=True)

def header(s, title, tag=None, color=DARK):
    rect(s, 0, 0, SW, Inches(1.15), color)
    rect(s, 0, Inches(1.15), SW, Pt(3), AMBER)
    rect(s, Inches(0.5), Inches(0.34), Inches(0.13), Inches(0.48), TEAL)
    tf = textbox(s, Inches(0.8), Inches(0.24), Inches(10.5), Inches(0.7))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    setrun(tf.paragraphs[0].add_run(), title, 26, WHITE, bold=True)
    if tag:
        tb = textbox(s, Inches(10.0), Inches(0.38), Inches(2.9), Inches(0.4))
        p = tb.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        setrun(p.add_run(), tag, 12, AMBER, bold=True)

def intro_line(s, text, top=1.42):
    tf = textbox(s, Inches(0.7), Inches(top), Inches(12.0), Inches(0.5))
    setrun(tf.paragraphs[0].add_run(), text, 15, TEAL, italic=True)

def bullets(s, items, top=2.05, left=0.7, width=12.0, height=4.7, gap=8):
    tf = textbox(s, Inches(left), Inches(top), Inches(width), Inches(height))
    first = True
    for it in items:
        lead = it.get("lead", ""); rest = it.get("rest", ""); lvl = it.get("level", 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.space_before = Pt(0)
        if lvl == 0:
            b = p.add_run(); setrun(b, "▸  ", 16, TEAL, bold=True)
            if lead:
                setrun(p.add_run(), lead, 16, INK, bold=True)
            if rest:
                setrun(p.add_run(), (" — " if lead else "") + rest, 15.5, GREY)
        else:
            p.level = 1
            b = p.add_run(); setrun(b, "–  ", 13, AMBER, bold=True)
            if lead:
                setrun(p.add_run(), lead, 13.5, INK, bold=True)
            if rest:
                setrun(p.add_run(), (" : " if lead else "") + rest, 13, GREY)
    return tf

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# =====================================================================
# SLIDE 1 — TITRE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(3.05), SW, Pt(3), TEAL)
rect(s, 0, Inches(4.35), SW, Pt(2), AMBER)
tf = textbox(s, Inches(1), Inches(1.15), Inches(11.3), Inches(0.5))
setrun(tf.paragraphs[0].add_run(),
       "École Nationale Supérieure d'Arts et Métiers — Meknès", 15, RGBColor(150,170,185))
tf = textbox(s, Inches(1), Inches(3.15), Inches(11.3), Inches(1.2))
setrun(tf.paragraphs[0].add_run(), "PrognoSense", 60, WHITE, bold=True)
tf = textbox(s, Inches(1), Inches(4.5), Inches(11.3), Inches(1.0))
setrun(tf.paragraphs[0].add_run(),
       "Plateforme Intelligente de Maintenance Prédictive Industrielle", 22, TEAL, bold=True)
tf = textbox(s, Inches(1), Inches(5.25), Inches(11.3), Inches(0.5))
setrun(tf.paragraphs[0].add_run(),
       "Analyse vibratoire  ·  Intelligence Artificielle  ·  Boucle fermée GMAO", 14,
       RGBColor(180,195,205), italic=True)
tf = textbox(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.9))
p = tf.paragraphs[0]
setrun(p.add_run(), "Réalisé par : ", 13, RGBColor(170,185,195))
setrun(p.add_run(), "Souleymane DIALLO", 13, WHITE, bold=True)
setrun(p.add_run(), "      Encadré par : ", 13, RGBColor(170,185,195))
setrun(p.add_run(), "Pr. Zaki SMAIL", 13, WHITE, bold=True)
p2 = tf.add_paragraph()
setrun(p2.add_run(),
       "Filière Intelligence Artificielle et Technologies des Données — Systèmes Industriels   ·   2025–2026",
       11, RGBColor(150,170,185))
notes(s, "Bonjour. Je vais vous présenter PrognoSense, une plateforme de maintenance "
         "prédictive industrielle que j'ai conçue et réalisée. Elle transforme l'analyse "
         "vibratoire en décisions de maintenance anticipées, de l'acquisition du signal "
         "jusqu'à l'ordre de travail. Je vais d'abord poser le contexte, puis dérouler "
         "l'architecture, les modèles, et chaque page de l'application.")

# =====================================================================
# SLIDE 2 — CONTEXTE & PROBLÉMATIQUE
# =====================================================================
s = slide(); header(s, "Contexte & problématique", "01")
intro_line(s, "Les arrêts non planifiés sont la source de pertes la plus lourde en industrie.")
bullets(s, [
 {"lead":"Maintenance corrective","rest":"réparer après la panne : arrêt subi, dégâts, coût maximal."},
 {"lead":"Maintenance préventive","rest":"intervalles fixes : gaspillage de pièces saines, pannes précoces non évitées."},
 {"lead":"Maintenance prédictive (CBM)","rest":"intervenir au bon moment, selon l'état réel — la plus efficiente."},
 {"lead":"Défi n°1","rest":"détecter la dégradation tôt, alors qu'elle est noyée dans le bruit."},
 {"lead":"Défi n°2","rest":"maîtriser le taux de fausses alarmes — sinon le système est débranché."},
 {"lead":"Défi n°3","rest":"rendre l'alerte actionnable et traçable (de l'alerte à l'intervention)."},
])
footer(s, 2)
notes(s, "L'enjeu central : plus de la moitié des défaillances de machines tournantes se "
         "manifestent à l'avance par les vibrations. Le vrai défi n'est pas seulement de "
         "détecter, mais de détecter TÔT, avec PEU de fausses alarmes, et de transformer "
         "l'alerte en action. C'est exactement ce que résout PrognoSense.")

# =====================================================================
# SLIDE 3 — OBJECTIFS & APPROCHE (OSA-CBM)
# =====================================================================
s = slide(); header(s, "Objectifs & approche", "02")
intro_line(s, "Couvrir toute la chaîne de surveillance conditionnelle (architecture OSA-CBM, ISO 13374).")
# chaine
chain = ["Acquisition","Traitement\ndu signal","Détection\nd'état","Indice\nde santé","Pronostic\n(RUL)","Action\n(GMAO)"]
x = Inches(0.7); y = Inches(2.05); w = Inches(1.85); h = Inches(0.95); gap = Inches(0.12)
for i, c in enumerate(chain):
    col = TEAL if i < 4 else AMBER
    sp = rect(s, x, y, w, h, col)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p.add_run(), c, 11.5, WHITE, bold=True)
    x = Emu(x + w + gap)
bullets(s, [
 {"lead":"Détecter","rest":"les dégradations et les anomalies, même inconnues."},
 {"lead":"Diagnostiquer","rest":"le défaut et sa sévérité selon la norme ISO 10816."},
 {"lead":"Pronostiquer","rest":"la durée de vie restante (RUL) pour planifier."},
 {"lead":"Agir & mesurer","rest":"ordres de travail GMAO + KPIs de fiabilité réels."},
 {"lead":"Faire confiance","rest":"explicabilité, audit, versioning — accessible à tous."},
], top=3.35)
footer(s, 3)
notes(s, "Je n'ai pas construit un simple détecteur, mais une chaîne complète conforme à "
         "l'architecture de référence OSA-CBM : acquisition, traitement, détection, santé, "
         "pronostic, action. Chaque maillon est implémenté. L'objectif final : passer "
         "d'un POC à une solution déployable, ouverte et explicable.")

# =====================================================================
# SLIDE 4 — ARCHITECTURE & TECHNOLOGIES
# =====================================================================
s = slide(); header(s, "Architecture & technologies", "03")
intro_line(s, "Une architecture en trois couches, modulaire et temps réel.")
layers = [("FRONTEND — React / Vite","7 pages · graphiques temps réel · WebSocket", GREEN),
          ("BACKEND — FastAPI","~20 routeurs REST + WebSocket · registre de modèles · pipeline unifié", TEAL),
          ("MODULES IA / SIGNAL","NumPy/SciPy · scikit-learn · XGBoost · PyTorch (auto-encodeur, CNN)", AMBER),
          ("PERSISTANCE","SQLAlchemy / SQLite (WAL) · TimescaleDB-ready", DARK)]
y = Inches(2.0)
for name, desc, col in layers:
    sp = rect(s, Inches(0.7), y, Inches(11.9), Inches(0.92), LIGHT)
    bar = rect(s, Inches(0.7), y, Inches(0.16), Inches(0.92), col)
    tf = textbox(s, Inches(1.05), y+Inches(0.08), Inches(11.4), Inches(0.8))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    setrun(p.add_run(), name + "   ", 14, col, bold=True)
    setrun(p.add_run(), desc, 12.5, GREY)
    y = Emu(y + Inches(1.06))
footer(s, 4)
notes(s, "L'architecture sépare clairement les responsabilités. Le backend FastAPI est "
         "asynchrone — indispensable pour le WebSocket temps réel. Un registre de modèles "
         "centralise toute l'inférence. Point clé : SQLite suffit pour la démo, mais le "
         "code bascule sur PostgreSQL/TimescaleDB en changeant une seule variable, sans "
         "rien réécrire.")

# =====================================================================
# SLIDE 5 — DATASETS & DÉFAUTS COUVERTS
# =====================================================================
s = slide(); header(s, "Datasets & défauts couverts", "04")
intro_line(s, "Cinq jeux de données de référence — plus de 66 000 fenêtres, multi-domaine et multi-tâche.")
rows = [["Dataset","Tâche","Fenêtres","Features"],
        ["VBL-VA001","Classification (4 cl.)","400","45"],
        ["CWRU","Classification (10 cl.)","12 627","15"],
        ["Mechanical Faults","Classification (4 cl.)","36 000","132"],
        ["CMAPSS FD001","Régression (RUL)","17 731","126"],
        ["MCC5-THU","Engrenages (CNN, 8 cl.)","—","—"]]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(2.0), Inches(6.6), Inches(3.0)).table
tbl.columns[0].width = Inches(2.2); tbl.columns[1].width = Inches(2.4)
tbl.columns[2].width = Inches(1.1); tbl.columns[3].width = Inches(0.9)
for j in range(4):
    cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
    r = cell.text_frame.paragraphs[0].add_run(); setrun(r, rows[0][j], 11.5, WHITE, bold=True)
for i in range(1, len(rows)):
    for j in range(4):
        cell = tbl.cell(i, j); cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
        r = cell.text_frame.paragraphs[0].add_run()
        setrun(r, rows[i][j], 10.5, INK, bold=(j==0))
tf = textbox(s, Inches(7.6), Inches(1.95), Inches(5.2), Inches(4.6))
setrun(tf.paragraphs[0].add_run(), "Modes de défaillance diagnostiqués", 14, AMBER, bold=True)
for t in ["Balourd, désalignement, jeu mécanique",
          "Roulement : bague ext./int./bille × 3 sévérités",
          "Engrenage : pitting, usure, dent cassée/fissurée/manquante, défauts mixtes",
          "Dégradation turbomoteur (pronostic RUL)",
          "Anomalie inconnue (auto-encodeur)"]:
    p = tf.add_paragraph(); p.space_before = Pt(6)
    setrun(p.add_run(), "▸ ", 12, TEAL, bold=True); setrun(p.add_run(), t, 12, GREY)
footer(s, 5)
notes(s, "La force du projet : là où la plupart des travaux traitent UN dataset et UNE "
         "tâche, je couvre quatre natures de machines et trois tâches — classification, "
         "régression RUL, et détection non supervisée. CWRU est exploité à 10 classes "
         "croisant la localisation ET la sévérité du défaut. C'est une vraie plateforme "
         "générique, pas un démonstrateur mono-cas.")

# =====================================================================
# SLIDE 6 — IA & BENCHMARK
# =====================================================================
s = slide(); header(s, "Intelligence artificielle & benchmark", "05")
intro_line(s, "Cinq modèles comparés — évaluation sans fuite de données (découpage par fichier source).")
rows = [["Dataset","Tâche","Meilleur modèle","Performance"],
        ["VBL-VA001","Classification (4 cl.)","XGBoost","100 %"],
        ["CWRU","Classification (10 cl.)","XGBoost","97,6 %"],
        ["Mechanical Faults","Classification (4 cl.)","MLP","89,6 %"],
        ["CMAPSS FD001","Régression RUL","XGBoost","MAE 9,61 · R² 0,90"]]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(2.05), Inches(11.9), Inches(2.7)).table
tbl.columns[0].width = Inches(2.7); tbl.columns[1].width = Inches(3.4)
tbl.columns[2].width = Inches(2.6); tbl.columns[3].width = Inches(3.2)
for j in range(4):
    cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = DARK
    r = cell.text_frame.paragraphs[0].add_run(); setrun(r, rows[0][j], 13, WHITE, bold=True)
for i in range(1, len(rows)):
    for j in range(4):
        cell = tbl.cell(i, j); cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
        r = cell.text_frame.paragraphs[0].add_run()
        col = GREEN if j == 3 else INK
        setrun(r, rows[i][j], 12.5, col, bold=(j>=2))
bullets(s, [
 {"lead":"Méthodologie rigoureuse","rest":"découpage par fichier source → aucune fuite de données ; score NASA asymétrique pour la RUL."},
 {"lead":"Aucun modèle universel","rest":"XGBoost gagne sur les roulements, le MLP sur Mechanical Faults → sélection par dataset, à chaud."},
], top=5.0)
footer(s, 6)
notes(s, "Ces chiffres sont réels et honnêtes : l'évaluation interdit toute fuite de "
         "données entre apprentissage et test. Sur CWRU, 97,6 % sur 10 classes. Sur le "
         "pronostic, une erreur moyenne de moins de 10 cycles avec un R² de 0,90. Le "
         "benchmark prouve qu'aucun modèle n'est supérieur partout : c'est pourquoi la "
         "plateforme permet de choisir le meilleur modèle par dataset, sans interruption.")

# =====================================================================
# SLIDE 7 — DÉTECTION D'ANOMALIE
# =====================================================================
s = slide(); header(s, "Détection d'anomalie & indice de santé", "06")
intro_line(s, "Trois approches complémentaires, dont une qui détecte l'inconnu.")
bullets(s, [
 {"lead":"Auto-encodeur (PyTorch)","rest":"entraîné uniquement sur des données saines → détecte les défauts JAMAIS vus."},
 {"lead":"Ensemble de 3 algorithmes","rest":"Isolation Forest + LOF + Elliptic Envelope, vote majoritaire → moins de fausses alertes."},
 {"lead":"Baseline propre à chaque machine","rest":"apprend le « normal » de CETTE machine (sain ≈ 5 %, dégradé ≈ 96 %), sans historique de panne."},
 {"lead":"Health Index (0–100 %)","rest":"fusion : 40 % anomalie + 40 % RUL + 20 % vibration → 4 seuils d'action."},
])
footer(s, 7)
notes(s, "C'est un point différenciant majeur : l'auto-encodeur, entraîné uniquement sur "
         "du sain, détecte des défauts qu'aucun classifieur supervisé ne connaît. "
         "L'ensemble à 3 votes réduit les fausses alarmes. Et le baseline par machine — "
         "le modèle d'Augury ou AspenTech — apprend le comportement sain propre à chaque "
         "équipement, sans avoir besoin d'une panne préalable.")

# =====================================================================
# Helper pour les slides PAGE
# =====================================================================
def page_slide(n, page_title, role, feats):
    s = slide(); header(s, page_title, f"PAGE {n:02d}", color=DARK)
    intro_line(s, "Rôle : " + role)
    bullets(s, feats, top=2.1)
    footer(s, n_global[0]); n_global[0]+=1
    return s

n_global = [8]

# SLIDE 8 — VUE GLOBALE
s = page_slide(1, "Page « Vue Globale » — supervision de la flotte",
 "voir l'état de toute la flotte d'un seul coup d'œil et prioriser.",
 [
  {"lead":"Cartes machines","rest":"jauge de santé, RUL estimée, tendance et statut par machine."},
  {"lead":"Codes couleur sémantiques","rest":"SAIN / SURVEILLANCE / ALERTE / CRITIQUE → décision immédiate."},
  {"lead":"KPIs de flotte","rest":"nombre de machines saines, en alerte, critiques, indice de santé moyen."},
  {"lead":"Sélection d'une machine","rest":"lance sa surveillance temps réel et un panneau détaillé (HI, RUL, anomalie, défaut)."},
 ])
notes(s, "La page d'accueil de l'ingénieur : la flotte est triée par criticité. En un "
         "regard, il sait quelles machines nécessitent son attention. Un clic bascule en "
         "surveillance temps réel.")

# SLIDE 9 — MONITORING
s = page_slide(2, "Page « Monitoring » — surveillance en direct",
 "suivre une machine en temps réel et voir la dégradation arriver.",
 [
  {"lead":"Forme d'onde & spectre temps réel","rest":"diffusés par WebSocket, sans rechargement."},
  {"lead":"Indicateurs vibratoires live","rest":"RMS, kurtosis, facteur de crête — détecter la dérive avant la panne."},
  {"lead":"Contrôles de simulation","rest":"play / pause / vitesse, sélection de la machine surveillée."},
  {"lead":"Reconnexion automatique","rest":"back-off exponentiel — robustesse de la liaison temps réel."},
 ])
notes(s, "Ici on suit une machine en direct : la courbe de santé et les indicateurs "
         "défilent. L'intérêt du prédictif est visible : on voit la tendance se dégrader "
         "progressivement, bien avant la panne.")

# SLIDE 10 — PRONOSTIC
s = page_slide(3, "Page « Pronostic & RUL » — anticiper la panne",
 "estimer la durée de vie restante et planifier l'intervention.",
 [
  {"lead":"Compte à rebours RUL","rest":"nombre de cycles restants avant défaillance, avec intervalle de confiance."},
  {"lead":"Trajectoire de dégradation","rest":"courbe de l'indice de santé avec seuils 70 / 40 / 20 %."},
  {"lead":"RUL prédite vs réelle","rest":"validation visuelle du pronostic (turbomoteurs CMAPSS)."},
  {"lead":"Valeur","rest":"planifier pièces de rechange et arrêts au bon moment, sans surprise."},
 ])
notes(s, "Le pronostic transforme la surveillance en planification : on sait combien de "
         "temps il reste, donc on commande les pièces et on programme l'arrêt au moment "
         "optimal — ni trop tôt, ni trop tard.")

# SLIDE 11 — IA LAB
s = page_slide(4, "Page « IA Lab » — maîtrise et confiance des modèles",
 "comparer, expliquer et gouverner les modèles d'IA.",
 [
  {"lead":"Benchmark & sélection à chaud","rest":"comparer les 5 modèles et changer le modèle actif sans interruption."},
  {"lead":"Fiabilité & dérive","rest":"score de drift (test de Kolmogorov-Smirnov) qui recommande un réentraînement."},
  {"lead":"Explicabilité & calibration","rest":"importance des indicateurs (SHAP), diagramme de fiabilité (ECE)."},
  {"lead":"Versions & rollback","rest":"historique des versions de modèles, retour arrière en un clic (MLOps)."},
 ])
notes(s, "Cette page incarne la confiance : on ne se contente pas de prédire, on explique "
         "sur quoi l'IA se base (SHAP), on surveille si elle reste valide (drift), et on "
         "peut revenir à une version antérieure si une mise à jour dégrade les "
         "performances. C'est de la gouvernance de modèles, du vrai MLOps.")

# SLIDE 12 — MAINTENANCE
s = page_slide(5, "Page « Maintenance » — de l'alerte à l'action",
 "le centre opérationnel : alertes, KPIs, ordres de travail et copilot.",
 [
  {"lead":"Copilot IA (RAG)","rest":"recommandations en langage naturel, citant les normes (sources tracées)."},
  {"lead":"Carnet de bord & KPIs réels","rest":"saisie d'interventions → MTBF, MTTR, disponibilité, ROI mesurés."},
  {"lead":"Efficacité prédictive","rest":"précision, rappel, lead-time + taux réel de fausses alarmes (validation expert)."},
  {"lead":"Ordres de travail automatiques","rest":"créés sur état critique, poussés vers la GMAO (SAP / Maximo)."},
  {"lead":"Réentraînement piloté","rest":"apprendre un nouveau défaut + diagnostic par spectre externe."},
 ])
notes(s, "C'est la page la plus riche, le cœur opérationnel. Trois idées fortes : "
         "(1) le copilot rend l'expertise accessible et cite ses sources ; "
         "(2) les KPIs de fiabilité sont MESURÉS sur les vraies interventions, pas "
         "inventés — un MTBF se constate, il ne se devine pas ; "
         "(3) la boucle est fermée : l'alerte devient un ordre de travail dans la GMAO.")

# SLIDE 13 — AUDIT / CONFIG / SECURITE
s = page_slide(6, "Audit Trail, Configuration & Sécurité",
 "traçabilité totale, paramétrage et accès sécurisé.",
 [
  {"lead":"Page Audit Trail","rest":"journal horodaté de chaque décision IA et alerte — conformité, « rien n'est une boîte noire »."},
  {"lead":"Page Configuration","rest":"réglage des seuils (anomalie, sévérité) et des paramètres de roulement."},
  {"lead":"Authentification & rôles","rest":"connexion JWT, rôles admin / utilisateur, mots de passe chiffrés (bcrypt)."},
  {"lead":"Validation des entrées","rest":"schémas Pydantic sur toute l'API (prévention des injections)."},
 ])
notes(s, "La traçabilité est essentielle en industrie : chaque alerte peut être justifiée "
         "a posteriori pour un audit. L'accès est sécurisé par rôles : seul un "
         "administrateur peut réentraîner un modèle ou faire un rollback.")

# =====================================================================
# SLIDE 14 — INDUSTRIALISATION & DIFFÉRENCIATION
# =====================================================================
s = slide(); header(s, "Industrialisation & différenciation", "07")
intro_line(s, "Ce qui fait passer le projet d'un démonstrateur à une solution déployable.")
bullets(s, [
 {"lead":"Connectivité terrain","rest":"ingestion universelle + connecteurs OPC-UA et MQTT — connectable au parc existant."},
 {"lead":"Diagnostic normalisé ISO 10816","rest":"vitesse en mm/s et zones A/B/C/D — défendable par un fiabiliste."},
 {"lead":"Boucle fermée GMAO","rest":"ordres de travail poussés vers SAP/Maximo + notifications e-mail."},
 {"lead":"Copilot ancré (RAG)","rest":"l'IA générative cite les normes et le guide des défauts."},
 {"lead":"Ouvert & explicable","rest":"open-source, sans lock-in — face aux suites propriétaires (SKF, Augury, GE) coûteuses et fermées."},
])
footer(s, 14)
notes(s, "Mon positionnement : je ne cherche pas à battre SKF ou GE en largeur. Je propose "
         "une alternative OUVERTE, EXPLICABLE et CONNECTABLE au parc existant via les "
         "standards OPC-UA/MQTT, qui parle le langage normalisé ISO, ferme la boucle "
         "jusqu'à la GMAO, et démocratise l'expertise par un copilot qui cite ses sources. "
         "C'est un message que peu d'acteurs portent aujourd'hui.")

# =====================================================================
# SLIDE 15 — RÉSULTATS, CONCLUSION & PERSPECTIVES
# =====================================================================
s = slide(); header(s, "Résultats, conclusion & perspectives", "08", color=DARK)
intro_line(s, "Une intégration verticale complète : du capteur jusqu'à l'ordre de travail.")
bullets(s, [
 {"lead":"Résultats clés","rest":"CWRU 97,6 % · Mechanical Faults 89,6 % · RUL MAE 9,61 / R² 0,90 · ISO validé · baseline 5 %/96 %."},
 {"lead":"Chaîne validée de bout en bout","rest":"ingestion OPC-UA, ordre de travail automatique, copilot RAG, KPIs réels."},
 {"lead":"Apports","rest":"rigueur (zéro fuite de données), confiance (audit, fausses alarmes), actionnabilité (GMAO)."},
], top=2.0, gap=10)
tf = textbox(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(2.2))
setrun(tf.paragraphs[0].add_run(), "Perspectives", 15, AMBER, bold=True)
for t in ["Validation sur données de terrain réelles + collecte run-to-failure",
          "Modèles physics-informed (hybrides physique + données), apprentissage fédéré multi-sites",
          "Durcissement cybersécurité IEC 62443, SSO/RBAC fin"]:
    p = tf.add_paragraph(); p.space_before = Pt(5)
    setrun(p.add_run(), "▸ ", 13, TEAL, bold=True); setrun(p.add_run(), t, 13, GREY)
p = tf.add_paragraph(); p.space_before = Pt(14)
setrun(p.add_run(), "Merci de votre attention — je suis à votre disposition pour vos questions.",
       15, TEAL, bold=True)
footer(s, 15)
notes(s, "En conclusion, PrognoSense est une intégration verticale complète : il ne se "
         "contente pas de prédire, il dit QUOI tombe en panne, POURQUOI le système le "
         "pense, DEPUIS QUAND il le sait, et QUELLE action mener — tout en restant "
         "vérifiable. Les résultats sont solides et honnêtes. Les perspectives ouvrent la "
         "voie vers une solution commercialisable. Merci, je suis prêt pour vos questions.")

import os
os.makedirs("presentation", exist_ok=True)
out = "presentation/PrognoSense_Soutenance.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
