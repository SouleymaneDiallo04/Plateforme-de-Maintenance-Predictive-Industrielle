# -*- coding: utf-8 -*-
"""Generateur du PowerPoint de soutenance PrognoSense — version 7 slides (10 min).

Centre sur les FONCTIONNALITES, langage ingenieur industriel, IA rendue accessible.
Charte sombre/cyan coherente avec l'application.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Palette (couleurs de l'application) ----
TEAL  = RGBColor(0, 150, 180)
DARK  = RGBColor(15, 34, 48)
AMBER = RGBColor(222, 150, 30)
GREEN = RGBColor(0, 160, 110)
RED   = RGBColor(210, 70, 70)
GREY  = RGBColor(80, 90, 100)
LIGHT = RGBColor(244, 248, 250)
WHITE = RGBColor(255, 255, 255)
INK   = RGBColor(28, 38, 48)
MUTE  = RGBColor(150, 170, 185)

N_TOTAL = 7

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color, line=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tf


def setrun(r, text, size, color, bold=False, italic=False, font="Calibri"):
    r.text = text; r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font


def footer(s, n):
    rect(s, Inches(0), Inches(7.15), SW, Pt(2), LIGHT)
    tf = textbox(s, Inches(0.5), Inches(7.04), Inches(9), Inches(0.4))
    setrun(tf.paragraphs[0].add_run(),
           "PrognoSense — Maintenance Prédictive Industrielle", 9, GREY)
    tf2 = textbox(s, Inches(11.3), Inches(7.04), Inches(1.6), Inches(0.4))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    setrun(p.add_run(), f"{n} / {N_TOTAL}", 9, TEAL, bold=True)


def header(s, title, tag=None):
    rect(s, 0, 0, SW, Inches(1.15), DARK)
    rect(s, 0, Inches(1.15), SW, Pt(3), AMBER)
    rect(s, Inches(0.5), Inches(0.34), Inches(0.13), Inches(0.48), TEAL)
    tf = textbox(s, Inches(0.8), Inches(0.2), Inches(10.3), Inches(0.78))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    setrun(tf.paragraphs[0].add_run(), title, 25, WHITE, bold=True)
    if tag:
        tb = textbox(s, Inches(10.4), Inches(0.4), Inches(2.5), Inches(0.4))
        p = tb.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        setrun(p.add_run(), tag, 12, AMBER, bold=True)


def intro_line(s, text, top=1.42):
    tf = textbox(s, Inches(0.7), Inches(top), Inches(12.0), Inches(0.5))
    setrun(tf.paragraphs[0].add_run(), text, 14.5, TEAL, italic=True)


def bullets(s, items, top=2.05, left=0.7, width=12.0, height=4.7, gap=9, lead_sz=15.5, rest_sz=14.5):
    tf = textbox(s, Inches(left), Inches(top), Inches(width), Inches(height))
    first = True
    for it in items:
        lead = it.get("lead", ""); rest = it.get("rest", "")
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.space_before = Pt(0)
        b = p.add_run(); setrun(b, "▸  ", 15, TEAL, bold=True)
        if lead:
            setrun(p.add_run(), lead, lead_sz, INK, bold=True)
        if rest:
            setrun(p.add_run(), (" — " if lead else "") + rest, rest_sz, GREY)
    return tf


def chip(s, l, t, w, h, text, fill, txt, sz=11.5, bold=True):
    sp = rect(s, l, t, w, h, fill, rounded=True)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p.add_run(), text, sz, txt, bold=bold)
    return sp


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# =====================================================================
# SLIDE 1 — TITRE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(3.0), SW, Pt(3), TEAL)
rect(s, 0, Inches(4.32), SW, Pt(2), AMBER)
tf = textbox(s, Inches(1), Inches(1.05), Inches(11.3), Inches(0.5))
setrun(tf.paragraphs[0].add_run(),
       "École Nationale Supérieure d'Arts et Métiers — Meknès", 15, MUTE)
tf = textbox(s, Inches(1), Inches(3.1), Inches(11.3), Inches(1.2))
setrun(tf.paragraphs[0].add_run(), "PrognoSense", 60, WHITE, bold=True)
tf = textbox(s, Inches(1), Inches(4.45), Inches(11.3), Inches(0.9))
setrun(tf.paragraphs[0].add_run(),
       "Plateforme de Maintenance Prédictive Industrielle", 23, TEAL, bold=True)
tf = textbox(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5))
setrun(tf.paragraphs[0].add_run(),
       "Surveiller · Diagnostiquer · Anticiper · Intervenir — du capteur à l'ordre de travail",
       14, MUTE, italic=True)
tf = textbox(s, Inches(1), Inches(6.25), Inches(11.3), Inches(0.9))
p = tf.paragraphs[0]
setrun(p.add_run(), "Réalisé par : ", 13, MUTE)
setrun(p.add_run(), "Souleymane DIALLO", 13, WHITE, bold=True)
setrun(p.add_run(), "        Encadré par : ", 13, MUTE)
setrun(p.add_run(), "Pr. Zaki SMAIL", 13, WHITE, bold=True)
p2 = tf.add_paragraph()
setrun(p2.add_run(),
       "Filière IA & Technologies des Données — Systèmes Industriels   ·   2025–2026", 11, MUTE)
notes(s, "Bonjour. Je vous présente PrognoSense, une plateforme de maintenance prédictive "
         "industrielle complète : elle surveille les machines, diagnostique le défaut précis, "
         "anticipe la panne, déclenche et suit l'intervention jusqu'à la pièce et au technicien. "
         "En 10 minutes je vais vous montrer ce qu'elle SAIT FAIRE, fonctionnalité par fonctionnalité.")

# =====================================================================
# SLIDE 2 — VUE D'ENSEMBLE : UNE PLATEFORME, 8 MODULES
# =====================================================================
s = slide(); header(s, "Une plateforme complète, 8 modules", "01")
intro_line(s, "Toute la surveillance conditionnelle (CBM), structurée comme la norme OSA-CBM / ISO 13374.")

# Chaine de valeur
chain = ["SURVEILLER", "DIAGNOSTIQUER", "ANTICIPER", "AGIR", "MESURER"]
sub   = ["état flotte", "organe défaillant", "durée de vie", "intervention", "fiabilité & ROI"]
x = Inches(0.7); y = Inches(2.05); w = Inches(2.28); h = Inches(1.0); gap = Inches(0.12)
for i, (c, sb) in enumerate(zip(chain, sub)):
    col = TEAL if i < 3 else AMBER
    sp = rect(s, x, y, w, h, col, rounded=True)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p.add_run(), c, 13, WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    setrun(p2.add_run(), sb, 9.5, WHITE)
    if i < len(chain) - 1:
        ar = textbox(s, Emu(x + w - Inches(0.02)), y + Inches(0.28), Inches(0.16), Inches(0.4))
        setrun(ar.paragraphs[0].add_run(), "›", 20, DARK, bold=True)
    x = Emu(x + w + gap)

# Les 8 espaces de travail
tf = textbox(s, Inches(0.7), Inches(3.45), Inches(12), Inches(0.4))
setrun(tf.paragraphs[0].add_run(), "8 espaces de travail dédiés :", 13, AMBER, bold=True)
modules = ["Vue Globale", "Monitoring", "Pronostic RUL", "IA Lab",
           "Maintenance", "Interventions", "Audit Trail", "Configuration"]
cw = Inches(2.95); ch = Inches(0.62); gx = Inches(0.13); gy = Inches(0.16)
x0 = Inches(0.7); y0 = Inches(3.95)
for i, m in enumerate(modules):
    r, c = divmod(i, 4)
    cx = Emu(x0 + c * (cw + gx)); cy = Emu(y0 + r * (ch + gy))
    chip(s, cx, cy, cw, ch, m, LIGHT, TEAL, sz=12.5)
tf = textbox(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.9))
setrun(tf.paragraphs[0].add_run(),
       "Une seule plateforme, du signal brut jusqu'à la décision d'intervention — "
       "pas un simple modèle, un système industriel complet.", 14, GREY, italic=True)
footer(s, 2)
notes(s, "PrognoSense couvre toute la chaîne de surveillance conditionnelle, dans l'ordre "
         "logique d'un fiabiliste : surveiller, diagnostiquer, anticiper, agir, mesurer. "
         "Concrètement c'est une application web de 8 modules. Dans les prochaines minutes "
         "je vais vous montrer les fonctionnalités de chacun. "
         "[Capture suggérée : la Vue Globale de la flotte.]")

# =====================================================================
# SLIDE 3 — ACQUERIR · SURVEILLER · DIAGNOSTIQUER
# =====================================================================
s = slide(); header(s, "Acquérir · Surveiller · Diagnostiquer", "02")
intro_line(s, "Capter le signal du parc existant et identifier le défaut — avec un diagnostic normalisé.")
bullets(s, [
 {"lead":"Raccordable au parc existant","rest":"connecteurs OPC-UA et MQTT + ingestion universelle : tout capteur ou automate alimente la plateforme."},
 {"lead":"Diagnostic à partir d'un simple spectre","rest":"on colle ou importe un spectre FFT (CSV/JSON) et le type de défaut est identifié automatiquement."},
 {"lead":"Surveillance temps réel","rest":"indice de santé par machine, forme d'onde, spectre, indicateurs (RMS, kurtosis, facteur de crête)."},
 {"lead":"Analyse spectrale experte","rest":"fréquences de roulement BPFO/BPFI/BSF/FTF + analyse d'enveloppe pour confirmer le défaut."},
 {"lead":"Sévérité normalisée ISO 10816","rest":"vitesse vibratoire en mm/s, classée en zones A/B/C/D — un diagnostic défendable, pas un indice maison."},
 {"lead":"Organe défaillant nommé précisément","rest":"« roulement bague interne », « désalignement d'accouplement »… — pas un libellé vague."},
])
footer(s, 3)
notes(s, "D'abord, d'où viennent les données : la plateforme se branche au parc via les "
         "standards industriels OPC-UA et MQTT — pas besoin de remplacer les capteurs. "
         "Fonctionnalité forte : on peut même diagnostiquer à partir d'un simple spectre FFT "
         "qu'on colle. La surveillance est temps réel, et surtout la sévérité est exprimée en "
         "zones ISO 10816 — le langage normalisé que vous connaissez. "
         "[Capture suggérée : page Monitoring ou diagnostic par spectre.]")

# =====================================================================
# SLIDE 4 — L'INTELLIGENCE : DECIDER, ANTICIPER, EXPLIQUER, S'AMELIORER
# =====================================================================
s = slide(); header(s, "L'intelligence : décider, anticiper, expliquer", "03")
intro_line(s, "Une IA qui diagnostique, pronostique, explique ses décisions — et s'améliore dans le temps.")
bullets(s, [
 {"lead":"Apprise sur des cas réels","rest":"5 jeux de référence (66 000+ cas) : roulements (CWRU), engrenages (MCC5-THU), balourd/désalignement/jeu, turbomoteur (CMAPSS)."},
 {"lead":"Plusieurs modèles, le meilleur retenu","rest":"arbres, Random Forest, XGBoost, réseau de neurones, CNN — le plus performant est choisi pour chaque cas."},
 {"lead":"Ce qu'elle produit","rest":"type de défaut + sévérité + pronostic de durée de vie restante (RUL) pour planifier l'intervention."},
 {"lead":"Détecter un défaut JAMAIS vu","rest":"elle apprend le « comportement sain » propre à chaque machine et signale toute déviation, même inconnue."},
 {"lead":"Elle explique ses décisions","rest":"elle montre quels indicateurs ont pesé (énergie BPFO, kurtosis…) — pas une boîte noire."},
 {"lead":"Elle s'améliore et reste fiable","rest":"réentraînement pour intégrer un nouveau défaut · retour arrière (rollback) en un clic · alerte de dérive."},
], top=1.95, gap=7, lead_sz=14.5, rest_sz=13.5)
# bandeau credibilite
rect(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5), LIGHT, rounded=True)
tf = textbox(s, Inches(0.95), Inches(6.5), Inches(11.5), Inches(0.4))
p = tf.paragraphs[0]
setrun(p.add_run(), "Fiabilité : ", 12.5, AMBER, bold=True)
setrun(p.add_run(), "jusqu'à 97 % de défauts correctement identifiés, validé sans tricher sur les données (aucune fuite entre apprentissage et test).", 12.5, INK)
footer(s, 4)
notes(s, "Le cerveau du système, expliqué simplement. Il a appris sur 5 jeux de données réels "
         "couvrant les défauts que vous connaissez : roulements, engrenages, balourd, "
         "désalignement, turbomachine. Trois fonctionnalités clés à retenir : (1) il détecte "
         "même un défaut JAMAIS rencontré, en apprenant le comportement sain de chaque machine ; "
         "(2) il EXPLIQUE pourquoi il décide — pas de boîte noire ; (3) il s'améliore : on peut "
         "lui apprendre un nouveau défaut, et revenir en arrière en un clic si besoin. "
         "Les résultats sont honnêtes : jusqu'à 97 % de bons diagnostics, sans fuite de données.")

# =====================================================================
# SLIDE 5 — DE L'ALERTE A L'INTERVENTION (boucle fermee)
# =====================================================================
s = slide(); header(s, "De l'alerte à l'intervention : la boucle fermée", "04")
intro_line(s, "L'alerte ne reste pas un voyant : elle devient une intervention affectée, équipée et tracée.")
bullets(s, [
 {"lead":"Ordre de travail automatique","rest":"créé sur état critique (priorité P1/P2), poussable vers la GMAO de l'usine (SAP PM, Maximo)."},
 {"lead":"Le bon technicien","rest":"affectation par compétence + certification ISO 18436 (l'analyste senior pour un cas grave) — l'admin choisit, le système classe les qualifiés."},
 {"lead":"La bonne pièce","rest":"vérification au magasin : la pièce requise est-elle en stock (quantité + emplacement) ou à commander ?"},
 {"lead":"Suivi de l'intervention","rest":"le technicien est notifié (e-mail), reçoit la mission détaillée, intervient, rédige son compte-rendu ; l'admin vérifie et clôture, le stock se décrémente."},
], top=1.95, gap=9)
# cycle de vie
states = ["CRÉÉ", "ASSIGNÉ", "EN COURS", "TERMINÉ", "CLÔTURÉ"]
x = Inches(0.9); y = Inches(5.55); w = Inches(2.05); h = Inches(0.62); gap = Inches(0.32)
for i, st in enumerate(states):
    col = GREEN if st == "CLÔTURÉ" else TEAL
    chip(s, x, y, w, h, st, col, WHITE, sz=12)
    if i < len(states) - 1:
        ar = textbox(s, Emu(x + w + Inches(0.02)), y + Inches(0.06), Inches(0.3), Inches(0.5))
        setrun(ar.paragraphs[0].add_run(), "→", 18, AMBER, bold=True)
    x = Emu(x + w + gap)
tf = textbox(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5))
setrun(tf.paragraphs[0].add_run(),
       "Le compte-rendu alimente automatiquement les KPIs réels et le taux de fausses alarmes.",
       12.5, GREY, italic=True)
footer(s, 5)
notes(s, "C'est le cœur opérationnel et l'innovation majeure. Quand une panne est prédite : "
         "un ordre de travail est créé automatiquement et peut partir dans votre GMAO. "
         "Ensuite on affecte LE BON technicien — selon sa compétence ET sa certification "
         "ISO 18436. On vérifie même si LA PIÈCE est disponible au magasin, ou s'il faut la "
         "commander. Le technicien est notifié, intervient, rédige son compte-rendu, l'admin "
         "vérifie et clôture. De la détection jusqu'à l'homme et à la pièce, la boucle est "
         "fermée et tracée. [Capture suggérée : page Interventions — dispatch + candidats + stock.]")

# =====================================================================
# SLIDE 6 — PILOTER, ASSISTER & TRACER
# =====================================================================
s = slide(); header(s, "Piloter la fiabilité, assister, tracer", "05")
intro_line(s, "Prouver la fiabilité par les chiffres, démocratiser l'expertise, garantir la traçabilité.")
bullets(s, [
 {"lead":"KPIs de fiabilité réels","rest":"MTBF, MTTR, disponibilité, ROI mesurés sur les interventions effectives — pas des chiffres inventés."},
 {"lead":"Taux de fausses alarmes maîtrisé","rest":"validé par l'expert sur le terrain (vrai/faux positif) — le critère de confiance n°1."},
 {"lead":"Assistant conversationnel (copilot)","rest":"le technicien pose une question en langage naturel ; il répond avec les sources normatives à l'appui."},
 {"lead":"Génération de rapport PDF","rest":"rapport de diagnostic / d'intervention exportable en un clic + notifications e-mail automatiques."},
 {"lead":"Journal d'audit","rest":"chaque décision et alerte horodatée → conformité et traçabilité ; rien n'est une boîte noire."},
 {"lead":"Page Configuration","rest":"seuils d'alarme, paramètres de roulement, préférences — la plateforme s'adapte à chaque site."},
], gap=8, lead_sz=15, rest_sz=14)
footer(s, 6)
notes(s, "Une fois qu'on intervient, il faut PILOTER et PROUVER. Les KPIs — MTBF, MTTR, "
         "disponibilité, ROI — sont mesurés sur les vraies interventions : un MTBF se "
         "constate, il ne s'invente pas. Le taux de fausses alarmes est validé par l'expert. "
         "À côté, des fonctionnalités qui font gagner du temps : un assistant qui répond aux "
         "questions en citant les normes, l'export de rapports PDF, les e-mails automatiques, "
         "le journal d'audit pour la conformité, et une page de configuration pour adapter les "
         "seuils à chaque site. [Capture suggérée : panneau KPIs + copilot, ou Audit Trail.]")

# =====================================================================
# SLIDE 7 — DEMO & CONCLUSION
# =====================================================================
s = slide(); header(s, "Démonstration & conclusion", "06")
intro_line(s, "Une intégration verticale complète : du capteur jusqu'à l'ordre de travail exécuté.")
bullets(s, [
 {"lead":"Démo express","rest":"flotte → une machine passe en alerte → page Interventions → affectation d'un technicien qualifié + pièce en stock."},
 {"lead":"Ce qui distingue PrognoSense","rest":"normalisé (ISO 10816/18436), ouvert (OPC-UA, MQTT, GMAO), expliqué (pas de boîte noire) et tracé (audit)."},
], top=1.95, gap=11, lead_sz=16, rest_sz=14.5)

tf = textbox(s, Inches(0.7), Inches(3.65), Inches(11.9), Inches(1.4))
setrun(tf.paragraphs[0].add_run(), "Perspectives", 15, AMBER, bold=True)
for t in ["Validation sur banc d'essai et données de terrain réelles",
          "Déploiement multi-sites (base time-series) et connecteurs GMAO natifs",
          "Application mobile pour les techniciens de terrain"]:
    p = tf.add_paragraph(); p.space_before = Pt(5)
    setrun(p.add_run(), "▸  ", 13, TEAL, bold=True); setrun(p.add_run(), t, 13.5, GREY)

# bandeau de cloture
rect(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.05), DARK, rounded=True)
tf = textbox(s, Inches(1.0), Inches(5.62), Inches(11.4), Inches(0.95))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
setrun(p.add_run(), "PrognoSense dit ", 16, WHITE)
setrun(p.add_run(), "quoi", 16, TEAL, bold=True)
setrun(p.add_run(), " tombe en panne, ", 16, WHITE)
setrun(p.add_run(), "pourquoi", 16, TEAL, bold=True)
setrun(p.add_run(), ", ", 16, WHITE)
setrun(p.add_run(), "quelle pièce", 16, TEAL, bold=True)
setrun(p.add_run(), " et ", 16, WHITE)
setrun(p.add_run(), "qui doit intervenir", 16, TEAL, bold=True)
setrun(p.add_run(), ".", 16, WHITE)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
setrun(p2.add_run(), "Merci de votre attention — je suis à votre disposition pour vos questions.",
       13.5, AMBER, bold=True)
footer(s, 7)
notes(s, "Pour finir, une démonstration rapide : une machine passe en alerte, et en quelques "
         "clics on affecte un technicien qualifié avec la pièce disponible. "
         "Ce qui distingue PrognoSense : c'est une intégration verticale complète, du capteur "
         "à l'ordre de travail — normalisée ISO, ouverte, expliquée et tracée. Là où les "
         "suites du marché sont fermées et coûteuses, la nôtre est ouverte et défendable. "
         "En une phrase : PrognoSense dit quoi tombe en panne, pourquoi, quelle pièce, et qui "
         "doit intervenir. Merci, je suis prêt pour vos questions.")

# ---------------------------------------------------------------------
import os
os.makedirs("presentation", exist_ok=True)
out = "presentation/PrognoSense_Soutenance_7slides.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
