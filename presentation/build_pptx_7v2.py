# -*- coding: utf-8 -*-
"""PowerPoint PrognoSense — 7 slides, design ESTHETIQUE (sombre, cartes, icones).

Style editorial/produit : fond sombre, tuiles a icones, chaines de process,
gros chiffres. Moins "rapport", plus "pitch produit". Coherent avec l'app.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Palette sombre (app) ----
BG    = RGBColor(13, 23, 33)
PANEL = RGBColor(22, 38, 52)
PANEL2= RGBColor(27, 46, 62)
CYAN  = RGBColor(0, 201, 224)
CYANd = RGBColor(0, 150, 180)
AMBER = RGBColor(240, 170, 50)
GREEN = RGBColor(45, 200, 140)
RED   = RGBColor(235, 95, 95)
WHITE = RGBColor(238, 244, 248)
MUTE  = RGBColor(140, 162, 178)
LINE  = RGBColor(42, 64, 82)

N_TOTAL = 7
EMO = "Segoe UI Emoji"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
MARGIN = 0.7


def slide():
    s = prs.slides.add_slide(BLANK)
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    sp.fill.solid(); sp.fill.fore_color.rgb = BG; sp.line.fill.background()
    sp.shadow.inherit = False
    return s


def rect(s, l, t, w, h, color, line=None, rounded=False, lw=1.0):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def tb(s, l, t, w, h, anchor=None):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def run(p, text, size, color, bold=False, italic=False, font="Calibri"):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return r


def para(tf, first=False, align=None, sb=0, sa=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align: p.alignment = align
    p.space_before = Pt(sb); p.space_after = Pt(sa)
    return p


def kicker(s, num, label):
    rect(s, MARGIN, 0.56, 0.34, 0.085, AMBER)
    tf = tb(s, MARGIN + 0.46, 0.42, 11.0, 0.4)
    p = para(tf, first=True)
    run(p, f"{num}", 12.5, AMBER, bold=True)
    run(p, "   " + label.upper(), 12.5, CYAN, bold=True)


def title(s, text, sz=30):
    tf = tb(s, MARGIN, 0.86, 12.0, 1.0)
    run(para(tf, first=True), text, sz, WHITE, bold=True)


def footer(s, n):
    rect(s, 0, 7.14, 13.333, 0.012, LINE)
    tf = tb(s, MARGIN, 7.02, 9, 0.35)
    run(para(tf, first=True), "PrognoSense  ·  Maintenance Prédictive Industrielle", 9, MUTE)
    tf2 = tb(s, 11.2, 7.02, 1.6, 0.35)
    run(para(tf2, first=True, align=PP_ALIGN.RIGHT), f"{n:02d} / {N_TOTAL:02d}", 9.5, CYAN, bold=True)


def icon_card(s, l, t, w, h, icon, ttl, desc, accent=CYAN, ttl_sz=13.5, desc_sz=11):
    rect(s, l, t, w, h, PANEL, line=LINE, rounded=True, lw=0.75)
    rect(s, l, t + 0.14, 0.07, h - 0.28, accent)            # barre d'accent gauche
    bd = rect(s, l + 0.22, t + 0.2, 0.62, 0.62, PANEL2, rounded=True)
    itf = bd.text_frame; itf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run(para(itf, first=True, align=PP_ALIGN.CENTER), icon, 22, accent, font=EMO)
    tf = tb(s, l + 1.02, t + 0.18, w - 1.2, 0.5)
    run(para(tf, first=True), ttl, ttl_sz, WHITE, bold=True)
    df = tb(s, l + 0.24, t + 0.86, w - 0.45, h - 1.0)
    run(para(df, first=True), desc, desc_sz, MUTE)


def stat(s, l, t, w, h, value, label, accent=CYAN):
    rect(s, l, t, w, h, PANEL, line=LINE, rounded=True, lw=0.75)
    tf = tb(s, l, t + 0.14, w, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), value, 27, accent, bold=True)
    lf = tb(s, l, t + h - 0.52, w, 0.4)
    run(para(lf, first=True, align=PP_ALIGN.CENTER), label, 10.5, MUTE, bold=True)


def flow_chips(s, items, x, y, max_w, accent=CYAN, fill=PANEL2, h=0.36, sz=11):
    cx, cy = x, y
    for it in items:
        w = 0.26 + 0.082 * len(it)
        if cx + w > x + max_w:
            cx = x; cy += h + 0.14
        rect(s, cx, cy, w, h, fill, line=LINE, rounded=True, lw=0.5)
        tf = tb(s, cx, cy, w, h, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, first=True, align=PP_ALIGN.CENTER), it, sz, accent, bold=True)
        cx += w + 0.14
    return cy + h


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# =====================================================================
# SLIDE 1 — HERO
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 0.10, CYAN)
tf = tb(s, MARGIN, 0.7, 11.9, 0.4)
run(para(tf, first=True), "ÉCOLE NATIONALE SUPÉRIEURE D'ARTS ET MÉTIERS — MEKNÈS   ·   2025–2026", 12, MUTE, bold=True)

tf = tb(s, MARGIN, 2.25, 11.9, 1.5)
run(para(tf, first=True), "PrognoSense", 66, WHITE, bold=True)
rect(s, MARGIN + 0.05, 3.62, 2.2, 0.06, CYAN)
tf = tb(s, MARGIN, 3.85, 11.9, 0.7)
run(para(tf, first=True), "Plateforme de Maintenance Prédictive Industrielle", 24, CYAN, bold=True)
tf = tb(s, MARGIN, 4.55, 11.9, 0.5)
run(para(tf, first=True),
    "Surveiller · Diagnostiquer · Anticiper · Intervenir — du capteur à l'ordre de travail.",
    14.5, MUTE, italic=True)

# pills thematiques
pills = ["Analyse vibratoire", "IA explicable", "ISO 10816 / 18436", "Boucle GMAO"]
x = MARGIN
for pz in pills:
    w = 0.4 + 0.092 * len(pz)
    rect(s, x, 5.35, w, 0.46, PANEL, line=LINE, rounded=True, lw=0.75)
    t = tb(s, x, 5.35, w, 0.46, anchor=MSO_ANCHOR.MIDDLE)
    run(para(t, first=True, align=PP_ALIGN.CENTER), pz, 12, CYAN, bold=True)
    x += w + 0.2

rect(s, MARGIN, 6.5, 11.93, 0.012, LINE)
tf = tb(s, MARGIN, 6.65, 11.9, 0.6)
p = para(tf, first=True)
run(p, "Réalisé par  ", 13, MUTE)
run(p, "Souleymane DIALLO", 14, WHITE, bold=True)
run(p, "          Encadré par  ", 13, MUTE)
run(p, "Pr. Zaki SMAIL", 14, WHITE, bold=True)
notes(s, "Bonjour. Je vous présente PrognoSense, une plateforme de maintenance prédictive "
         "industrielle complète. En 10 minutes, je vais vous montrer ce qu'elle SAIT FAIRE — "
         "fonctionnalité par fonctionnalité — de la surveillance jusqu'à l'intervention, le "
         "technicien et la pièce.")

# =====================================================================
# SLIDE 2 — VUE D'ENSEMBLE
# =====================================================================
s = slide(); kicker(s, "01", "La plateforme")
title(s, "Une seule plateforme, du capteur à la décision")

# chaine de valeur
chain = [("SURVEILLER", "état flotte", CYAN), ("DIAGNOSTIQUER", "organe défaillant", CYAN),
         ("ANTICIPER", "durée de vie", CYAN), ("AGIR", "intervention", AMBER),
         ("MESURER", "fiabilité · ROI", AMBER)]
x = MARGIN; w = 2.18; h = 1.0; gap = 0.27; y = 1.95
for i, (c, sb, col) in enumerate(chain):
    rect(s, x, y, w, h, PANEL, line=col, rounded=True, lw=1.2)
    tf = tb(s, x, y + 0.16, w, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, align=PP_ALIGN.CENTER), c, 12.5, WHITE, bold=True)
    run(para(tf, align=PP_ALIGN.CENTER, sb=1), sb, 9.5, col)
    if i < len(chain) - 1:
        a = tb(s, x + w - 0.04, y + 0.28, 0.34, 0.45)
        run(para(a, first=True, align=PP_ALIGN.CENTER), "›", 22, AMBER, bold=True)
    x += w + gap

# 8 modules
tf = tb(s, MARGIN, 3.35, 11.9, 0.35)
run(para(tf, first=True), "8 ESPACES DE TRAVAIL DÉDIÉS", 12, AMBER, bold=True)
mods = [("🗂️","Vue Globale"), ("📈","Monitoring"), ("⏳","Pronostic RUL"), ("🧠","IA Lab"),
        ("🛠️","Maintenance"), ("👷","Interventions"), ("🛡️","Audit Trail"), ("⚙️","Configuration")]
cw = 2.84; ch = 0.92; gx = 0.18; gy = 0.18; x0 = MARGIN; y0 = 3.78
for i, (ic, nm) in enumerate(mods):
    r, c = divmod(i, 4)
    cx = x0 + c * (cw + gx); cy = y0 + r * (ch + gy)
    rect(s, cx, cy, cw, ch, PANEL, line=LINE, rounded=True, lw=0.75)
    itf = tb(s, cx + 0.18, cy, 0.7, ch, anchor=MSO_ANCHOR.MIDDLE)
    run(para(itf, first=True, align=PP_ALIGN.CENTER), ic, 21, CYAN, font=EMO)
    ntf = tb(s, cx + 0.82, cy, cw - 0.9, ch, anchor=MSO_ANCHOR.MIDDLE)
    run(para(ntf, first=True), nm, 13, WHITE, bold=True)
tf = tb(s, MARGIN, 5.95, 11.9, 0.7)
run(para(tf, first=True),
    "Toute la surveillance conditionnelle (CBM / ISO 13374) dans un seul outil — "
    "pas un modèle isolé, un système industriel complet.", 13.5, MUTE, italic=True)
footer(s, 2)
notes(s, "PrognoSense suit l'ordre logique d'un fiabiliste : surveiller, diagnostiquer, "
         "anticiper, agir, mesurer. Concrètement, c'est une application de 8 modules que je "
         "vais parcourir. L'idée à retenir : ce n'est pas un détecteur de plus, c'est la "
         "chaîne complète, du signal jusqu'à la décision d'intervention.")

# =====================================================================
# SLIDE 3 — ACQUERIR / SURVEILLER / DIAGNOSTIQUER (6 tuiles)
# =====================================================================
s = slide(); kicker(s, "02", "Surveillance & diagnostic")
title(s, "Capter le signal, voir le défaut, le nommer")
cards3 = [
 ("📡", "Raccordable au parc", "Connecteurs OPC-UA & MQTT + ingestion universelle : tout capteur ou automate alimente la plateforme.", CYAN),
 ("📊", "Diagnostic par spectre", "On colle/importe un spectre FFT (CSV/JSON) → le type de défaut est identifié automatiquement.", CYAN),
 ("⏱️", "Surveillance temps réel", "Indice de santé, forme d'onde, spectre, indicateurs (RMS, kurtosis, facteur de crête).", CYAN),
 ("🎯", "Analyse spectrale experte", "Fréquences de roulement BPFO/BPFI/BSF/FTF + analyse d'enveloppe pour confirmer.", AMBER),
 ("📏", "Sévérité ISO 10816", "Vitesse vibratoire en mm/s, zones A/B/C/D — un diagnostic normalisé, défendable.", AMBER),
 ("🔩", "Organe défaillant nommé", "« roulement bague interne », « désalignement d'accouplement »… — jamais un libellé vague.", AMBER),
]
cw = (11.93 - 2 * 0.3) / 3; ch = 1.95; x0 = MARGIN; y0 = 2.05
for i, (ic, tt, de, ac) in enumerate(cards3):
    r, c = divmod(i, 3)
    icon_card(s, x0 + c * (cw + 0.3), y0 + r * (ch + 0.28), cw, ch, ic, tt, de, ac)
footer(s, 3)
notes(s, "D'abord, d'où viennent les données : la plateforme se branche au parc existant via "
         "OPC-UA et MQTT — pas besoin de changer les capteurs. On peut même diagnostiquer à "
         "partir d'un simple spectre collé. La surveillance est temps réel, et surtout la "
         "sévérité est exprimée en zones ISO 10816 — votre langage. Et le défaut est nommé "
         "précisément, pas en termes vagues.")

# =====================================================================
# SLIDE 4 — INTELLIGENCE (composition mixte)
# =====================================================================
s = slide(); kicker(s, "03", "Intelligence embarquée")
title(s, "Diagnostiquer, anticiper, expliquer — et s'améliorer")

# Panneau gauche : LE MOTEUR
LX, LW = MARGIN, 5.55
rect(s, LX, 2.0, LW, 4.05, PANEL, line=LINE, rounded=True, lw=0.75)
tf = tb(s, LX + 0.28, 2.18, LW - 0.5, 0.4)
run(para(tf, first=True), "LE MOTEUR", 12.5, AMBER, bold=True)
tf = tb(s, LX + 0.28, 2.55, LW - 0.55, 0.55)
run(para(tf, first=True), "Appris sur 5 jeux de données réels — 66 000+ cas", 12.5, WHITE, bold=True)
yy = flow_chips(s, ["CWRU", "MCC5-THU", "Mechanical Faults", "VBL", "CMAPSS"],
                LX + 0.28, 3.08, LW - 0.55, accent=CYAN)
tf = tb(s, LX + 0.28, yy + 0.12, LW - 0.55, 0.4)
run(para(tf, first=True), "Modèles comparés, le meilleur retenu par cas :", 12, MUTE)
yy = flow_chips(s, ["XGBoost", "Random Forest", "MLP", "CNN", "Huber"],
                LX + 0.28, yy + 0.5, LW - 0.55, accent=GREEN, fill=PANEL2)
tf = tb(s, LX + 0.28, yy + 0.18, LW - 0.55, 0.7)
p = para(tf, first=True)
run(p, "Produit : ", 12.5, AMBER, bold=True)
run(p, "type de défaut + sévérité + durée de vie restante (RUL).", 12.5, WHITE)

# Colonne droite : 3 capacites
RX, RW = 6.55, 6.08
caps = [
 ("🔎", "Détecter un défaut JAMAIS vu", "Apprend le « comportement sain » de chaque machine et signale toute déviation, même inconnue.", CYAN),
 ("💡", "Expliquer ses décisions", "Montre quels indicateurs ont pesé (énergie BPFO, kurtosis…) — pas une boîte noire.", AMBER),
 ("🔄", "S'améliorer dans le temps", "Réentraînement pour un nouveau défaut · retour arrière (rollback) en un clic · alerte de dérive.", GREEN),
]
ch = 1.18; gy = 0.16; y0 = 2.0
for i, (ic, tt, de, ac) in enumerate(caps):
    icon_card(s, RX, y0 + i * (ch + gy), RW, ch, ic, tt, de, ac, ttl_sz=13, desc_sz=10.5)

# bandeau credibilite
rect(s, MARGIN, 6.3, 11.93, 0.62, PANEL2, line=CYAN, rounded=True, lw=1.0)
tf = tb(s, MARGIN + 0.3, 6.3, 11.4, 0.62, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True)
run(p, "✓  Fiabilité  ", 13.5, GREEN, bold=True, font=EMO)
run(p, "jusqu'à 97 % de défauts correctement identifiés — validé sans tricher (aucune fuite de données entre apprentissage et test).", 13, WHITE)
footer(s, 4)
notes(s, "Le cerveau, expliqué simplement. Il a appris sur 5 jeux de données réels couvrant "
         "vos défauts : roulements, engrenages, balourd, désalignement, turbomachine. Trois "
         "choses à retenir : il détecte même un défaut jamais vu ; il EXPLIQUE pourquoi il "
         "décide ; et il s'améliore — on lui apprend un nouveau défaut, et on revient en "
         "arrière en un clic si besoin. Résultats honnêtes : jusqu'à 97 %, sans fuite de données.")

# =====================================================================
# SLIDE 5 — BOUCLE FERMEE (process + cartes)
# =====================================================================
s = slide(); kicker(s, "04", "De l'alerte à l'intervention")
title(s, "La boucle fermée : la bonne personne, la bonne pièce")

# process horizontal
steps = [("⚠️","Alerte", RED), ("📋","Ordre de\ntravail", AMBER), ("👷","Technicien\nqualifié", CYAN),
         ("📦","Pièce en\nstock", CYAN), ("✅","Clôture\nvérifiée", GREEN)]
x = MARGIN; w = 2.0; h = 1.25; gap = 0.36; y = 1.95
for i, (ic, lb, col) in enumerate(steps):
    rect(s, x, y, w, h, PANEL, line=col, rounded=True, lw=1.1)
    itf = tb(s, x, y + 0.16, w, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    run(para(itf, first=True, align=PP_ALIGN.CENTER), ic, 24, col, font=EMO)
    ltf = tb(s, x, y + 0.66, w, 0.55)
    for k, ln in enumerate(lb.split("\n")):
        run(para(ltf, first=(k == 0), align=PP_ALIGN.CENTER), ln, 11.5, WHITE, bold=True)
    if i < len(steps) - 1:
        a = tb(s, x + w - 0.02, y + 0.42, 0.4, 0.5)
        run(para(a, first=True, align=PP_ALIGN.CENTER), "→", 20, AMBER, bold=True)
    x += w + gap

# 3 cartes detaillees
cards5 = [
 ("🎓", "Affectation par compétence", "Compétence + certification ISO 18436 : le senior pour un cas grave. L'admin choisit, le système classe les qualifiés.", CYAN),
 ("📦", "Vérification au magasin", "La pièce requise est-elle en stock (quantité + emplacement) ou à commander ? Décrément automatique à la clôture.", AMBER),
 ("📝", "Suivi & compte-rendu", "Technicien notifié par e-mail, mission détaillée, compte-rendu, vérification admin — qui alimente seul les KPIs.", GREEN),
]
cw = (11.93 - 2 * 0.3) / 3; ch = 1.85; y0 = 3.5
for i, (ic, tt, de, ac) in enumerate(cards5):
    icon_card(s, MARGIN + i * (cw + 0.3), y0, cw, ch, ic, tt, de, ac, desc_sz=10.5)
tf = tb(s, MARGIN, 5.55, 11.9, 0.4)
footer(s, 5)
notes(s, "C'est le cœur opérationnel et l'innovation majeure. Une panne prédite déclenche un "
         "ordre de travail automatique, poussable dans votre GMAO. On affecte LE BON "
         "technicien — selon compétence ET certification ISO 18436. On vérifie si LA PIÈCE "
         "est au magasin ou à commander. Le technicien est notifié, intervient, rédige son "
         "compte-rendu, l'admin clôture. De la détection jusqu'à l'homme et la pièce, la "
         "boucle est fermée et tracée.")

# =====================================================================
# SLIDE 6 — PILOTAGE & CONFIANCE (stats + cartes)
# =====================================================================
s = slide(); kicker(s, "05", "Pilotage & confiance")
title(s, "Prouver la fiabilité, assister, tracer")

# stats KPI mesurees
sw = (11.93 - 3 * 0.28) / 4; y = 1.95; h = 1.15
for i, (v, lb, col) in enumerate([("MTBF", "temps entre pannes", CYAN), ("MTTR", "temps de réparation", CYAN),
                                   ("DISPO.", "disponibilité", GREEN), ("ROI", "arrêts & coûts évités", AMBER)]):
    stat(s, MARGIN + i * (sw + 0.28), y, sw, h, v, lb, col)
tf = tb(s, MARGIN, 3.18, 11.9, 0.35)
run(para(tf, first=True), "KPIS MESURÉS SUR LES INTERVENTIONS RÉELLES — PAS DES CHIFFRES INVENTÉS", 11.5, MUTE, bold=True)

cards6 = [
 ("🤖", "Assistant conversationnel", "Le technicien pose une question en langage naturel ; réponse avec les sources normatives à l'appui.", CYAN),
 ("📄", "Rapports & alertes", "Génération de rapport PDF en un clic + notifications e-mail automatiques sur événement critique.", CYAN),
 ("🛡️", "Journal d'audit", "Chaque décision et alerte horodatée → conformité & traçabilité. Rien n'est une boîte noire.", AMBER),
 ("🎚️", "Taux de fausses alarmes", "Validé par l'expert (vrai/faux positif) — le critère de confiance n°1 d'un acheteur PdM.", AMBER),
 ("⚙️", "Configuration par site", "Seuils d'alarme, paramètres de roulement, préférences : la plateforme s'adapte à chaque usine.", GREEN),
 ("📈", "Efficacité prédictive", "Précision, rappel, délai d'anticipation (lead-time) entre alertes et pannes réelles.", GREEN),
]
cw = (11.93 - 2 * 0.3) / 3; ch = 1.32; x0 = MARGIN; y0 = 3.55
for i, (ic, tt, de, ac) in enumerate(cards6):
    r, c = divmod(i, 3)
    icon_card(s, x0 + c * (cw + 0.3), y0 + r * (ch + 0.16), cw, ch, ic, tt, de, ac, ttl_sz=12.5, desc_sz=10)
footer(s, 6)
notes(s, "Une fois qu'on intervient, il faut PILOTER et PROUVER. MTBF, MTTR, disponibilité, "
         "ROI sont mesurés sur les vraies interventions — un MTBF se constate, il ne s'invente "
         "pas. À côté, des fonctionnalités qui font gagner du temps : un assistant qui répond "
         "en citant les normes, l'export PDF, les e-mails, le journal d'audit pour la "
         "conformité, le taux de fausses alarmes validé, et une configuration par site.")

# =====================================================================
# SLIDE 7 — CONCLUSION
# =====================================================================
s = slide(); kicker(s, "06", "Conclusion")
title(s, "Une intégration verticale complète")

diff = [("📐","Normalisé", "ISO 10816 / 18436", CYAN),
        ("🔌","Ouvert", "OPC-UA · MQTT · GMAO", CYAN),
        ("💡","Expliqué", "pas de boîte noire", AMBER),
        ("🛡️","Tracé", "audit & conformité", GREEN)]
cw = (11.93 - 3 * 0.28) / 4; y0 = 2.0; h = 1.35
for i, (ic, tt, de, ac) in enumerate(diff):
    l = MARGIN + i * (cw + 0.28)
    rect(s, l, y0, cw, h, PANEL, line=LINE, rounded=True, lw=0.75)
    itf = tb(s, l, y0 + 0.16, cw, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    run(para(itf, first=True, align=PP_ALIGN.CENTER), ic, 23, ac, font=EMO)
    ttf = tb(s, l, y0 + 0.66, cw, 0.4)
    run(para(ttf, first=True, align=PP_ALIGN.CENTER), tt, 14, WHITE, bold=True)
    dtf = tb(s, l, y0 + 1.0, cw, 0.35)
    run(para(dtf, first=True, align=PP_ALIGN.CENTER), de, 10.5, MUTE)

# perspectives
tf = tb(s, MARGIN, 3.7, 11.9, 1.2)
run(para(tf, first=True), "PERSPECTIVES", 12, AMBER, bold=True)
for t in ["Validation sur banc d'essai et données de terrain réelles",
          "Déploiement multi-sites (base time-series) · connecteurs GMAO natifs · application mobile terrain"]:
    p = para(tf, sb=4)
    run(p, "▸  ", 13, CYAN, bold=True); run(p, t, 13, MUTE)

# bandeau de cloture
rect(s, MARGIN, 5.2, 11.93, 1.25, PANEL, line=CYAN, rounded=True, lw=1.2)
tf = tb(s, MARGIN + 0.4, 5.32, 11.1, 1.05, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True)
run(p, "PrognoSense dit ", 17, WHITE)
run(p, "quoi", 17, CYAN, bold=True)
run(p, " tombe en panne, ", 17, WHITE)
run(p, "pourquoi", 17, CYAN, bold=True)
run(p, ", ", 17, WHITE)
run(p, "quelle pièce", 17, CYAN, bold=True)
run(p, ", et ", 17, WHITE)
run(p, "qui doit intervenir", 17, CYAN, bold=True)
run(p, ".", 17, WHITE)
p2 = para(tf, sb=6)
run(p2, "Merci de votre attention — je suis à votre disposition pour vos questions.", 13, AMBER, bold=True)
footer(s, 7)
notes(s, "Pour finir : PrognoSense est une intégration verticale complète, du capteur à "
         "l'ordre de travail — normalisée ISO, ouverte, expliquée et tracée. Là où les suites "
         "du marché sont fermées et coûteuses, la nôtre est ouverte et défendable. En une "
         "phrase : elle dit quoi tombe en panne, pourquoi, quelle pièce, et qui doit "
         "intervenir. Merci, je suis prêt pour vos questions.")

# ---------------------------------------------------------------------
import os
os.makedirs("presentation", exist_ok=True)
out = "presentation/PrognoSense_Soutenance_v2_design.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
